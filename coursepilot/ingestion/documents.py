"""Local preprocessing for supported course material files."""

import hashlib
from collections.abc import Sequence
from pathlib import Path
from typing import Protocol

import fitz  # type: ignore[import-untyped]
from pptx import Presentation
from pydantic import BaseModel, ConfigDict, Field

from coursepilot.models import MaterialType


class UnsupportedFileType(ValueError):
    """Raised when an uploaded file is not PDF or PPTX."""


class FileTooLarge(ValueError):
    """Raised when an uploaded file exceeds the configured limit."""


class EmptyExtraction(ValueError):
    """Raised when a supported document contains no extractable text."""


class PageContent(BaseModel):
    model_config = ConfigDict(frozen=True)

    page_number: int = Field(ge=1)
    text: str


class MaterialParser(Protocol):
    def supports(self, file_name: str) -> bool: ...

    def parse(self, path: Path) -> list[PageContent]: ...


class UploadValidator:
    def __init__(self, *, max_upload_bytes: int) -> None:
        if max_upload_bytes <= 0:
            raise ValueError("max_upload_bytes must be positive")
        self._max_upload_bytes = max_upload_bytes

    def validate(self, path: Path) -> MaterialType:
        extension = path.suffix.casefold()
        try:
            material_type = {".pdf": MaterialType.PDF, ".pptx": MaterialType.PPTX}[extension]
        except KeyError as error:
            raise UnsupportedFileType(f"unsupported file extension: {extension}") from error
        if path.stat().st_size > self._max_upload_bytes:
            raise FileTooLarge(f"file exceeds {self._max_upload_bytes} bytes")
        return material_type


def calculate_file_hash(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as source:
        while chunk := source.read(1024 * 1024):
            digest.update(chunk)
    return digest.hexdigest()


class PdfParser:
    def supports(self, file_name: str) -> bool:
        return Path(file_name).suffix.casefold() == ".pdf"

    def parse(self, path: Path) -> list[PageContent]:
        with fitz.open(path) as document:
            pages = [
                PageContent(page_number=index + 1, text=page.get_text("text").strip())
                for index, page in enumerate(document)
            ]
        _ensure_text(pages, path)
        return pages


class PptxParser:
    def supports(self, file_name: str) -> bool:
        return Path(file_name).suffix.casefold() == ".pptx"

    def parse(self, path: Path) -> list[PageContent]:
        presentation = Presentation(str(path))
        pages: list[PageContent] = []
        for index, slide in enumerate(presentation.slides, start=1):
            text_parts = [
                shape.text.strip()
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            pages.append(PageContent(page_number=index, text="\n".join(text_parts)))
        _ensure_text(pages, path)
        return pages


class MarkdownRenderer:
    def render(self, title: str, pages: Sequence[PageContent], material_type: MaterialType) -> str:
        label = "PDF" if material_type is MaterialType.PDF else "PPT"
        sections = [f"# {title.strip()}"]
        sections.extend(f"## {label} 第 {page.page_number} 页\n\n{page.text}" for page in pages)
        return "\n\n".join(sections).rstrip() + "\n"


def _ensure_text(pages: Sequence[PageContent], path: Path) -> None:
    if not any(page.text.strip() for page in pages):
        raise EmptyExtraction(f"no extractable text in {path.name}")
