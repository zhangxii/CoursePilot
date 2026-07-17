from pathlib import Path

import fitz
import pytest
from pptx import Presentation
from pptx.util import Inches

from coursepilot.ingestion import (
    EmptyExtraction,
    FileTooLarge,
    MarkdownRenderer,
    PdfParser,
    PptxParser,
    UnsupportedFileType,
    UploadValidator,
    calculate_file_hash,
)
from coursepilot.models import MaterialType


def create_pdf(path: Path, texts: list[str]) -> None:
    document = fitz.open()
    for text in texts:
        page = document.new_page()
        if text:
            page.insert_text((72, 72), text)
    document.save(path)
    document.close()


def create_pptx(path: Path, texts: list[str]) -> None:
    presentation = Presentation()
    for text in texts:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        box.text = text
    presentation.save(path)


def test_upload_validator_accepts_pdf_pptx_and_rejects_invalid_files(tmp_path: Path) -> None:
    pdf = tmp_path / "course.PDF"
    pdf.write_bytes(b"pdf")
    pptx = tmp_path / "course.pptx"
    pptx.write_bytes(b"pptx")
    text = tmp_path / "course.txt"
    text.write_text("text", encoding="utf-8")
    oversized = tmp_path / "large.pdf"
    oversized.write_bytes(b"x" * 11)
    validator = UploadValidator(max_upload_bytes=10)

    assert validator.validate(pdf) is MaterialType.PDF
    assert validator.validate(pptx) is MaterialType.PPTX
    with pytest.raises(UnsupportedFileType):
        validator.validate(text)
    with pytest.raises(FileTooLarge):
        validator.validate(oversized)


def test_file_hash_is_stable_and_changes_with_content(tmp_path: Path) -> None:
    first = tmp_path / "first.pdf"
    second = tmp_path / "second.pdf"
    changed = tmp_path / "changed.pdf"
    first.write_bytes(b"same")
    second.write_bytes(b"same")
    changed.write_bytes(b"different")

    assert calculate_file_hash(first) == calculate_file_hash(second)
    assert calculate_file_hash(first) != calculate_file_hash(changed)


def test_pdf_parser_extracts_pages_in_order_and_rejects_empty_document(tmp_path: Path) -> None:
    source = tmp_path / "course.pdf"
    create_pdf(source, ["Page one", "Page two"])

    pages = PdfParser().parse(source)

    assert [(page.page_number, page.text) for page in pages] == [
        (1, "Page one"),
        (2, "Page two"),
    ]

    empty = tmp_path / "empty.pdf"
    create_pdf(empty, [""])
    with pytest.raises(EmptyExtraction):
        PdfParser().parse(empty)


def test_pptx_parser_extracts_slides_in_order_and_rejects_empty_deck(tmp_path: Path) -> None:
    source = tmp_path / "course.pptx"
    create_pptx(source, ["Slide one", "Slide two"])

    pages = PptxParser().parse(source)

    assert [(page.page_number, page.text) for page in pages] == [
        (1, "Slide one"),
        (2, "Slide two"),
    ]

    empty = tmp_path / "empty.pptx"
    create_pptx(empty, [""])
    with pytest.raises(EmptyExtraction):
        PptxParser().parse(empty)


def test_markdown_renderer_preserves_source_type_and_page_numbers(tmp_path: Path) -> None:
    pdf = tmp_path / "course.pdf"
    pptx = tmp_path / "course.pptx"
    create_pdf(pdf, ["PDF content"])
    create_pptx(pptx, ["PPT content"])
    renderer = MarkdownRenderer()

    pdf_markdown = renderer.render("PDF Course", PdfParser().parse(pdf), MaterialType.PDF)
    pptx_markdown = renderer.render("PPT Course", PptxParser().parse(pptx), MaterialType.PPTX)

    assert "# PDF Course" in pdf_markdown
    assert "## PDF 第 1 页" in pdf_markdown
    assert "PDF content" in pdf_markdown
    assert "# PPT Course" in pptx_markdown
    assert "## PPT 第 1 页" in pptx_markdown
    assert "PPT content" in pptx_markdown
